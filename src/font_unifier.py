import sys
import os
import logging

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFileDialog, QMessageBox, QFrame,
    QComboBox, QCompleter, QProgressBar, QStyle
)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QEvent
from PyQt6.QtGui import QFont, QFontDatabase
from docx import Document
from docx.oxml.ns import qn as docx_qn
from openpyxl import load_workbook
from pptx import Presentation
from pptx.oxml.ns import qn as pptx_qn
from lxml import etree


logger = logging.getLogger(__name__)


# --- Font helpers (handle East Asian / complex scripts) ---

def _set_docx_run_font(run, font_name):
    """Word: set latin + eastAsia + complex-script fonts on a run.

    The theme-reference attributes (asciiTheme/eastAsiaTheme/cstheme/...) are
    removed too: when they coexist with an explicit name, some viewers fall
    back to the theme font instead of the explicit one (the same class of
    issue as Excel's <scheme>).
    """
    rfonts = run._element.get_or_add_rPr().get_or_add_rFonts()
    rfonts.set(docx_qn('w:ascii'), font_name)
    rfonts.set(docx_qn('w:hAnsi'), font_name)
    rfonts.set(docx_qn('w:eastAsia'), font_name)
    rfonts.set(docx_qn('w:cs'), font_name)
    for theme_attr in ('w:asciiTheme', 'w:hAnsiTheme',
                       'w:eastAsiaTheme', 'w:cstheme'):
        rfonts.attrib.pop(docx_qn(theme_attr), None)


def _set_pptx_run_font(run, font_name):
    """PowerPoint: set latin + eastAsian + complex-script typefaces on a run."""
    # run.font.name への代入は rPr を確実に生成し、a:latin も設定する
    run.font.name = font_name
    rPr = run.font._rPr
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        elem = rPr.find(pptx_qn(tag))
        if elem is None:
            elem = etree.SubElement(rPr, pptx_qn(tag))
        elem.set('typeface', font_name)


def _set_pptx_text_frame_fonts(text_frame, font_name):
    """Apply the pptx font helper to every run in a text frame."""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            _set_pptx_run_font(run, font_name)


# --- Core Logic for Font Changing ---

def _set_docx_font(paragraphs, font_name):
    """Apply the docx font helper to every run across the given paragraphs."""
    for para in paragraphs:
        for run in para.runs:
            _set_docx_run_font(run, font_name)


def _process_docx_container(container, font_name):
    """Process paragraphs + tables of a body/header/footer/cell.

    A cell may itself contain nested tables, so this recurses naturally via
    _process_docx_table -> _process_docx_container.
    """
    _set_docx_font(container.paragraphs, font_name)
    for table in getattr(container, 'tables', ()):
        _process_docx_table(table, font_name)


def _process_docx_table(table, font_name):
    for row in table.rows:
        for cell in row.cells:
            _process_docx_container(cell, font_name)


def change_word_font(path, font_name):
    """Changes the font for all text in a .docx file.

    Covers body paragraphs/tables (incl. nested tables) and per-section
    headers/footers (default, first-page, even-page). Drawing text boxes
    (<w:txbxContent>) are not covered — known limitation.
    """
    doc = Document(path)
    _process_docx_container(doc, font_name)
    for section in doc.sections:
        for part in (section.header, section.footer,
                     section.first_page_header, section.first_page_footer,
                     section.even_page_header, section.even_page_footer):
            _process_docx_container(part, font_name)
    return doc


def _replace_all_fonts(workbook, font_name):
    """Replace the name of every font definition in the workbook.

    All cells, named styles and the default (Normal) style reference these
    shared font definitions, so updating them in place covers the whole
    workbook — including unstyled/empty cells that otherwise keep the old
    default font across sheets. Other font attributes (size, bold, ...) are
    preserved.

    The ``scheme`` attribute (minor/major) is also cleared: when present,
    Excel ignores the explicit <name> and renders text with the theme font
    (e.g. the East-Asian minor font), so previously unstyled cells would
    keep showing the old font even after the name was changed.

    注意: workbook._fonts は openpyxl 3.x の非公開 API。依存バージョンは
    requirements.txt で openpyxl==3.1.5 に固定済み。アップグレード時は再検証が必要。
    """
    for font in workbook._fonts:
        font.name = font_name
        font.scheme = None


def change_excel_font(path, font_name):
    """Changes the font for all cells in a .xlsx file, preserving other style."""
    workbook = load_workbook(path)
    _replace_all_fonts(workbook, font_name)
    return workbook


def _process_chart_fonts(chart, font_name):
    """Set fonts on a chart's title and axis titles.

    Wrapped defensively: a malformed chart must never crash the whole file.
    python-pptx exposes category_axis/value_axis (not x_axis/y_axis). Per-point
    data-label fonts are not reliably settable, so they are skipped.
    """
    try:
        if chart.has_title:
            _set_pptx_text_frame_fonts(chart.chart_title.text_frame, font_name)
        for axis_attr in ('category_axis', 'value_axis', 'series_axis'):
            axis = getattr(chart, axis_attr, None)
            if axis is not None and getattr(axis, 'has_title', False):
                try:
                    _set_pptx_text_frame_fonts(
                        axis.axis_title.text_frame, font_name)
                except Exception:
                    logger.debug("chart axis font failed", exc_info=True)
    except Exception:
        logger.debug("chart font failed", exc_info=True)


def _process_ppt_shape(shape, font_name):
    """Recursively processes text in a shape, including nested groups."""
    if getattr(shape, 'has_text_frame', False):
        _set_pptx_text_frame_fonts(shape.text_frame, font_name)
    if getattr(shape, 'has_table', False):
        for row in shape.table.rows:
            for cell in row.cells:
                _set_pptx_text_frame_fonts(cell.text_frame, font_name)
    if getattr(shape, 'has_chart', False):
        _process_chart_fonts(shape.chart, font_name)
    if getattr(shape, 'has_group', False):
        for sub_shape in shape.shapes:
            _process_ppt_shape(sub_shape, font_name)


def change_ppt_font(path, font_name):
    """Changes the font for all text in a .pptx file.

    Covers text frames, tables, charts (title/axes) and nested groups.
    """
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            _process_ppt_shape(shape, font_name)
    return prs


# Extension -> handler (case-insensitive dispatch in process_office_file)
_FONT_CHANGERS = {
    ".docx": change_word_font,
    ".xlsx": change_excel_font,
    ".pptx": change_ppt_font,
}


def process_office_file(path, font_name):
    """Process a single Office file and save the modified copy.

    Returns the output path. Raises ValueError on unsupported extensions.
    Case-insensitive on the extension.
    """
    root, ext = os.path.splitext(path)
    output_path = f"{root}_modified{ext}"

    changer = _FONT_CHANGERS.get(ext.lower())
    if changer is None:
        raise ValueError(f"Unsupported file type: {ext}")
    changer(path, font_name).save(output_path)
    return output_path


# --- Background worker (keeps the GUI responsive on large files) ---

class FontProcessingWorker(QThread):
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, path, font_name):
        super().__init__()
        self._path = path
        self._font_name = font_name

    def run(self):
        try:
            output_path = process_office_file(self._path, self._font_name)
            self.finished.emit(output_path)
        except Exception as e:
            self.error.emit(str(e))


# --- UI Styling ---

ACCENT = "#2563EB"
ACCENT_HOVER = "#1D4ED8"
ACCENT_PRESSED = "#1E40AF"
ACCENT_DISABLED = "#94A3B8"
BG = "#F1F5F9"
CARD = "#FFFFFF"
FILE_CARD_BG = "#F8FAFC"
GHOST_HOVER = "#EFF6FF"
TEXT = "#0F172A"
MUTED = "#64748B"
BORDER = "#E2E8F0"

# type/badge -> (background, foreground)
BADGE_COLORS = {
    "docx": ("#DBEAFE", "#1E40AF"),
    "xlsx": ("#DCFCE7", "#166534"),
    "pptx": ("#FFEDD5", "#9A3412"),
}
# status kind -> (background, foreground)
STATUS_COLORS = {
    "success": ("#DCFCE7", "#166534"),
    "error": ("#FEE2E2", "#991B1B"),
    "info": ("#DBEAFE", "#1E40AF"),
}


def _color_qss(selector, prop, colors):
    """Render QSS rules: selector[prop="key"] { background; color }."""
    return "\n".join(
        f'{selector}[{prop}="{key}"] {{ background: {bg}; color: {fg}; }}'
        for key, (bg, fg) in colors.items())


APP_QSS = f"""
QMainWindow, QWidget#central {{ background: {BG}; }}

QFrame#cardFrame {{
    background: {CARD};
    border: 1px solid {BORDER};
    border-radius: 12px;
}}
QFrame#fileCard {{
    background: {FILE_CARD_BG};
    border: 1px solid {BORDER};
    border-radius: 8px;
}}

QLabel#titleLabel {{
    color: {TEXT};
    font-size: 18pt;
    font-weight: bold;
}}
QLabel#subtitleLabel {{
    color: {MUTED};
    font-size: 10pt;
}}

QLabel#badge {{
    border-radius: 10px;
    padding: 2px 10px;
    font-size: 9pt;
    font-weight: bold;
}}
{_color_qss('QLabel#badge', 'type', BADGE_COLORS)}

QPushButton#primary {{
    background: {ACCENT};
    color: #FFFFFF;
    border: none;
    border-radius: 8px;
    padding: 10px 28px;
    font-size: 11pt;
    font-weight: bold;
}}
QPushButton#primary:hover {{ background: {ACCENT_HOVER}; }}
QPushButton#primary:pressed {{ background: {ACCENT_PRESSED}; }}
QPushButton#primary:disabled {{ background: {ACCENT_DISABLED}; color: {FILE_CARD_BG}; }}

QPushButton#ghost {{
    background: transparent;
    color: {ACCENT};
    border: 1px solid {ACCENT};
    border-radius: 6px;
    padding: 6px 14px;
}}
QPushButton#ghost:hover {{ background: {GHOST_HOVER}; }}

QLineEdit, QComboBox {{
    background: {CARD};
    border: 1px solid {BORDER};
    border-radius: 6px;
    padding: 6px 8px;
    color: {TEXT};
    selection-background-color: {ACCENT};
    selection-color: #FFFFFF;
}}
QLineEdit:focus, QComboBox:focus {{ border: 1px solid {ACCENT}; }}
QLineEdit:disabled {{ color: {MUTED}; }}
QComboBox::drop-down {{ border: none; width: 20px; }}
QComboBox QAbstractItemView {{
    background: {CARD};
    border: 1px solid {BORDER};
    selection-background-color: {ACCENT};
    selection-color: #FFFFFF;
    outline: none;
}}

QProgressBar {{
    background: {BORDER};
    border: none;
    border-radius: 4px;
    max-height: 8px;
}}
QProgressBar::chunk {{ background: {ACCENT}; border-radius: 4px; }}

QLabel#statusLabel {{
    border-radius: 6px;
    padding: 6px 12px;
    font-size: 10pt;
}}
{_color_qss('QLabel#statusLabel', 'kind', STATUS_COLORS)}
"""


# --- GUI Application ---

class FontUnifierApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Font Unifier")
        self.resize(620, 520)
        self.setMinimumSize(560, 480)

        self.file_path = ""
        self.font_name = "Meiryo UI"
        self._worker = None

        central = QWidget()
        central.setObjectName("central")
        self.setCentralWidget(central)
        outer = QVBoxLayout(central)
        outer.setContentsMargins(24, 24, 24, 24)

        card = QFrame()
        card.setObjectName("cardFrame")
        outer.addWidget(card)
        layout = QVBoxLayout(card)
        layout.setContentsMargins(28, 28, 28, 28)
        layout.setSpacing(16)

        # Header
        title = QLabel("Font Unifier")
        title.setObjectName("titleLabel")
        layout.addWidget(title)
        subtitle = QLabel("批量统一 Office 文档字体")
        subtitle.setObjectName("subtitleLabel")
        layout.addWidget(subtitle)

        # Type badges
        badge_row = QHBoxLayout()
        badge_row.setSpacing(8)
        for ext in ("docx", "xlsx", "pptx"):
            badge = QLabel(ext.upper())
            badge.setObjectName("badge")
            badge.setProperty("type", ext)
            badge_row.addWidget(badge)
        badge_row.addStretch()
        layout.addLayout(badge_row)

        # File card
        file_card = QFrame()
        file_card.setObjectName("fileCard")
        file_inner = QVBoxLayout(file_card)
        file_inner.setContentsMargins(14, 12, 14, 12)
        file_inner.setSpacing(8)

        file_head = QHBoxLayout()
        file_head.addWidget(self._muted_label("选择文件"))
        file_head.addStretch()
        browse_button = QPushButton("Browse…")
        browse_button.setObjectName("ghost")
        browse_icon = self.style().standardIcon(
            QStyle.StandardPixmap.SP_DialogOpenButton)
        browse_button.setIcon(browse_icon)
        browse_button.setCursor(Qt.CursorShape.PointingHandCursor)
        browse_button.clicked.connect(self.browse_file)
        file_head.addWidget(browse_button)
        file_inner.addLayout(file_head)

        file_row = QHBoxLayout()
        file_icon = QLabel()
        file_icon.setPixmap(browse_icon.pixmap(20, 20))
        file_row.addWidget(file_icon)
        self.file_entry = QLineEdit()
        self.file_entry.setReadOnly(True)
        self.file_entry.setPlaceholderText("未选择文件")
        file_row.addWidget(self.file_entry, 1)
        file_inner.addLayout(file_row)
        layout.addWidget(file_card)

        # Font row
        font_row = QHBoxLayout()
        font_label = self._muted_label("目标字体")
        font_label.setFixedWidth(80)
        font_row.addWidget(font_label)
        self.font_entry = QComboBox()
        self.font_entry.setEditable(True)
        font_families = QFontDatabase.families()
        self.font_entry.addItems(font_families)
        default = self.font_name if self.font_name in font_families else (
            font_families[0] if font_families else "")
        self.font_entry.setCurrentText(default)
        # 入力時に前方一致で候補をポップアップ表示（大小区別なし）
        completer = QCompleter(font_families, self)
        completer.setCaseSensitivity(Qt.CaseSensitivity.CaseInsensitive)
        completer.setFilterMode(Qt.MatchFlag.MatchStartsWith)
        completer.setCompletionMode(QCompleter.CompletionMode.PopupCompletion)
        self.font_entry.setCompleter(completer)
        # 入力欄のクリックでドロップダウンを開く（Excel のフォント選択と同様の挙動）
        self.font_entry.lineEdit().installEventFilter(self)
        font_row.addWidget(self.font_entry, 1)
        layout.addLayout(font_row)

        # Progress (busy spinner, hidden until processing)
        self.progress = QProgressBar()
        self.progress.setTextVisible(False)
        self.progress.setRange(0, 0)
        self.progress.setVisible(False)
        layout.addWidget(self.progress)

        # Action
        self.start_button = QPushButton("Start Processing")
        self.start_button.setObjectName("primary")
        self.start_button.setCursor(Qt.CursorShape.PointingHandCursor)
        self.start_button.clicked.connect(self.process_file)
        layout.addWidget(self.start_button,
                         alignment=Qt.AlignmentFlag.AlignCenter)

        # Status
        self.status_label = QLabel("")
        self.status_label.setObjectName("statusLabel")
        self.status_label.setProperty("kind", "idle")
        self.status_label.setVisible(False)
        layout.addWidget(self.status_label,
                         alignment=Qt.AlignmentFlag.AlignCenter)

        layout.addStretch()

    def _muted_label(self, text):
        label = QLabel(text)
        label.setStyleSheet(f"color: {MUTED}; font-weight: bold;")
        return label

    def _finish_processing(self):
        self.progress.setVisible(False)
        self.start_button.setEnabled(True)

    def closeEvent(self, event):
        # 処理中にウィンドウを閉じた場合、スレッド終了を待ってから破棄する
        if self._worker is not None and self._worker.isRunning():
            self._worker.wait(5000)
        event.accept()

    def eventFilter(self, obj, event):
        # フォント入力欄のクリックで候補リストを開く（入力時は前方可動で絞り込まれる）
        if obj is self.font_entry.lineEdit() and \
                event.type() == QEvent.Type.MouseButtonPress:
            completer = self.font_entry.completer()
            completer.setCompletionPrefix("")
            completer.complete()
        return super().eventFilter(obj, event)

    def _set_status(self, text, kind):
        self.status_label.setText(text)
        self.status_label.setProperty("kind", kind)
        self.status_label.style().unpolish(self.status_label)
        self.status_label.style().polish(self.status_label)
        self.status_label.setVisible(kind != "idle")

    def browse_file(self):
        file_dialog = QFileDialog(self)
        file_dialog.setNameFilters([
            "Office Files (*.docx *.xlsx *.pptx)",
            "Word Documents (*.docx)",
            "Excel Workbooks (*.xlsx)",
            "PowerPoint Presentations (*.pptx)",
            "All files (*.*)"
        ])
        if file_dialog.exec():
            selected_files = file_dialog.selectedFiles()
            if selected_files:
                self.file_path = selected_files[0]
                self.file_entry.setText(self.file_path)
                self._set_status("", "idle")

    def process_file(self):
        path = self.file_path
        font = self.font_entry.currentText()

        if not path:
            QMessageBox.critical(self, "Error", "Please select a file first.")
            return
        if not font:
            QMessageBox.critical(self, "Error",
                                 "Please enter a target font name.")
            return

        self._set_status("Processing...", "info")
        self.progress.setVisible(True)
        self.start_button.setEnabled(False)

        self._worker = FontProcessingWorker(path, font)
        self._worker.finished.connect(self._on_processing_finished)
        self._worker.error.connect(self._on_processing_error)
        self._worker.start()

    def _on_processing_finished(self, output_path):
        self._finish_processing()
        self._set_status(f"Success! Saved to {output_path}", "success")
        QMessageBox.information(
            self, "Success",
            "File processed successfully and saved as: " + output_path)

    def _on_processing_error(self, message):
        self._finish_processing()
        self._set_status("An error occurred.", "error")
        QMessageBox.critical(
            self, "Error",
            "An error occurred during processing: " + message)


if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setStyleSheet(APP_QSS)
    app.setFont(QFont("Segoe UI", 10))
    window = FontUnifierApp()
    window.show()
    sys.exit(app.exec())
