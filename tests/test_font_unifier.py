import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from docx import Document  # noqa: E402
from openpyxl import Workbook, load_workbook  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

import font_unifier  # noqa: E402

TARGET_FONT = "Arial"


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def _make_docx(path):
    """段落(複数 run)と表格を含む docx を生成する"""
    doc = Document()
    # 複数 run を持つ段落
    para = doc.add_paragraph()
    para.add_run("Hello ")
    para.add_run("World")
    # 表格
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "C"
    table.cell(1, 1).text = "D"
    doc.save(str(path))
    return str(path)


def test_change_word_font_updates_all_runs(tmp_path):
    path = _make_docx(tmp_path / "in.docx")
    doc = font_unifier.change_word_font(path, TARGET_FONT)

    # 段落内の全 run を検証
    for para in doc.paragraphs:
        for run in para.runs:
            assert run.font.name == TARGET_FONT

    # 表格内の全 run を検証
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        assert run.font.name == TARGET_FONT


def test_change_word_font_empty_doc_no_error(tmp_path):
    """段落・表格を持たない空の docx でも例外を投げない"""
    doc = Document()
    # デフォルトで 1 つの空段落が付くが run は無い -> そのまま保存
    empty_path = str(tmp_path / "empty.docx")
    doc.save(empty_path)

    result = font_unifier.change_word_font(empty_path, TARGET_FONT)
    assert result is not None


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------

def _make_xlsx(path):
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Score"
    ws["A2"] = "Alice"
    ws["B2"] = 90
    ws.title = "Data"
    wb.save(str(path))
    return str(path)


def test_change_excel_font_updates_non_empty_cells(tmp_path):
    path = _make_xlsx(tmp_path / "in.xlsx")
    wb = font_unifier.change_excel_font(path, TARGET_FONT)

    for sheetname in wb.sheetnames:
        ws = wb[sheetname]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    assert cell.font.name == TARGET_FONT, (
                        f"cell {cell.coordinate} font={cell.font.name}"
                    )


def test_change_excel_font_persists_after_save(tmp_path):
    """保存して再読込してもフォントが維持されることを検証"""
    path = _make_xlsx(tmp_path / "in.xlsx")
    wb = font_unifier.change_excel_font(path, TARGET_FONT)
    out = str(tmp_path / "out.xlsx")
    wb.save(out)

    wb2 = load_workbook(out)
    ws = wb2.active
    assert ws["A1"].font.name == TARGET_FONT
    assert ws["B2"].font.name == TARGET_FONT


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def _make_pptx(path):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.text = "Title"
    p = tf.paragraphs[0]
    p.add_run().text = " Extra"
    prs.save(str(path))
    return str(path)


def test_change_ppt_font_updates_textbox_runs(tmp_path):
    path = _make_pptx(tmp_path / "in.pptx")
    prs = font_unifier.change_ppt_font(path, TARGET_FONT)

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        assert run.font.name == TARGET_FONT


def test_change_ppt_font_persists_after_save(tmp_path):
    path = _make_pptx(tmp_path / "in.pptx")
    prs = font_unifier.change_ppt_font(path, TARGET_FONT)
    out = str(tmp_path / "out.pptx")
    prs.save(out)

    prs2 = Presentation(out)
    slide = prs2.slides[0]
    shape = slide.shapes[0]
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            assert run.font.name == TARGET_FONT


# ---------------------------------------------------------------------------
# 追加: S1 Excel スタイル保持 / S2 東アジアフォント / S3 チャート / S4 拡張子
# ---------------------------------------------------------------------------

def test_change_excel_font_preserves_other_style(tmp_path):
    """S1: フォント名のみ変更し、サイズ/太字などは保持する"""
    from openpyxl.styles import Font
    path = str(tmp_path / "styled.xlsx")
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "X"
    ws["A1"].font = Font(name="Calibri", size=20, bold=True)
    wb.save(path)

    wb2 = font_unifier.change_excel_font(path, TARGET_FONT)
    f = wb2.active["A1"].font
    assert f.name == TARGET_FONT
    assert f.size == 20
    assert f.bold is True


def test_change_word_font_sets_east_asia(tmp_path):
    """S2: Word の run に eastAsia/ascii 属性が設定される"""
    from docx.oxml.ns import qn
    path = _make_docx(tmp_path / "in.docx")
    doc = font_unifier.change_word_font(path, TARGET_FONT)
    run = doc.paragraphs[0].runs[0]
    rfonts = run._element.rPr.rFonts
    assert rfonts.get(qn('w:ascii')) == TARGET_FONT
    assert rfonts.get(qn('w:eastAsia')) == TARGET_FONT
    assert rfonts.get(qn('w:cs')) == TARGET_FONT


def test_change_ppt_font_sets_east_asian_typeface(tmp_path):
    """S2: PowerPoint の run に <a:ea> typeface が設定される"""
    from pptx.oxml.ns import qn
    path = _make_pptx(tmp_path / "in.pptx")
    prs = font_unifier.change_ppt_font(path, TARGET_FONT)
    run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    rPr = run.font._rPr
    assert rPr.find(qn('a:latin')).get('typeface') == TARGET_FONT
    assert rPr.find(qn('a:ea')).get('typeface') == TARGET_FONT
    assert rPr.find(qn('a:cs')).get('typeface') == TARGET_FONT


def _make_pptx_with_chart(path, with_data_labels=False):
    """タイトル付きチャート(必要ならデータラベル有効化)を含む pptx を生成"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ['A', 'B']
    cd.add_series('S1', (1, 2))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1), Inches(4), Inches(3), cd)
    gframe.chart.has_title = True
    gframe.chart.chart_title.text_frame.text = "ChartTitle"
    if with_data_labels:
        gframe.chart.series[0].data_labels.show_value = True
    prs.save(str(path))
    return str(path)


def test_change_ppt_font_chart_title_and_no_crash(tmp_path):
    """S3: チャートを含む pptx がクラッシュせず、タイトル run のフォントが更新される"""
    from pptx.oxml.ns import qn
    path = _make_pptx_with_chart(tmp_path / "chart.pptx")
    prs = font_unifier.change_ppt_font(path, TARGET_FONT)
    chart = prs.slides[0].shapes[0].chart
    run = chart.chart_title.text_frame.paragraphs[0].runs[0]
    assert run.font.name == TARGET_FONT
    assert run.font._rPr.find(qn('a:ea')).get('typeface') == TARGET_FONT


def test_change_ppt_font_chart_with_data_labels_no_crash(tmp_path):
    """S3: データラベル有効なチャートでもクラッシュしない(元実装はここで TypeError)"""
    path = _make_pptx_with_chart(
        tmp_path / "chart_dl.pptx", with_data_labels=True)
    prs = font_unifier.change_ppt_font(path, TARGET_FONT)
    assert prs is not None


def test_change_ppt_font_table_cells(tmp_path):
    """PPT テーブルセル内の run のフォントが更新される"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    gtbl = slide.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    gtbl.cell(0, 0).text = "X"
    gtbl.cell(1, 1).text = "Y"
    path = str(tmp_path / "table.pptx")
    prs.save(path)

    prs2 = font_unifier.change_ppt_font(path, TARGET_FONT)
    tbl = prs2.slides[0].shapes[0].table
    found = 0
    for row in tbl.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        assert run.font.name == TARGET_FONT
                        found += 1
    assert found >= 2


def test_process_office_file_case_insensitive_extension(tmp_path):
    """S4: 大文字拡張子 .PPTX を受け付ける"""
    path = _make_pptx(tmp_path / "in.pptx")
    upper_path = str(tmp_path / "UP.PPTX")
    os.rename(path, upper_path)

    out = font_unifier.process_office_file(upper_path, TARGET_FONT)
    assert os.path.exists(out)
    assert out.endswith("_modified.PPTX")


def test_process_office_file_unsupported_extension(tmp_path):
    """未対応拡張子は ValueError"""
    path = str(tmp_path / "file.txt")
    with open(path, "w", encoding="utf-8") as fh:
        fh.write("hi")
    try:
        font_unifier.process_office_file(path, TARGET_FONT)
    except ValueError:
        return
    raise AssertionError("ValueError was expected for .txt")
