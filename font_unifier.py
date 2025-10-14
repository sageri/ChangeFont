import sys
import os
from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFileDialog, QMessageBox, QFrame
)
from PyQt6.QtCore import Qt
from docx import Document
from docx.shared import Pt
from openpyxl import load_workbook
from openpyxl.styles import Font
from pptx import Presentation

# --- Core Logic for Font Changing ---

def change_word_font(path, new_font_name):
    """Changes the font for all text in a .docx file."""
    doc = Document(path)
    for para in doc.paragraphs:
        for run in para.runs:
            run.font.name = new_font_name
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for run in para.runs:
                        run.font.name = new_font_name
    
    return doc

def change_excel_font(path, new_font_name):
    """Changes the font for all cells in a .xlsx file."""
    workbook = load_workbook(path)
    new_font = Font(name=new_font_name)
    
    for sheetname in workbook.sheetnames:
        worksheet = workbook[sheetname]
        for row in worksheet.iter_rows():
            for cell in row:
                # We apply the font to all cells, even empty ones, to ensure consistency
                # when data is added later.
                cell.font = new_font
                
    return workbook

def change_ppt_font(path, new_font_name):
    """Changes the font for all text in a .pptx file."""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = new_font_name
    return prs

# --- GUI Application ---

class FontUnifierApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Font Unifier")
        self.setGeometry(100, 100, 500, 250)

        self.file_path = ""
        self.font_name = "Meiryo UI"

        # Central widget
        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        # Main layout
        layout = QVBoxLayout(central_widget)

        # File Selection Frame
        file_frame = QFrame()
        file_layout = QHBoxLayout(file_frame)
        file_layout.setContentsMargins(10, 10, 10, 10)

        file_label = QLabel("File:")
        file_label.setFixedWidth(60)
        file_layout.addWidget(file_label)

        self.file_entry = QLineEdit()
        self.file_entry.setReadOnly(True)
        file_layout.addWidget(self.file_entry)

        browse_button = QPushButton("Browse...")
        browse_button.clicked.connect(self.browse_file)
        file_layout.addWidget(browse_button)

        layout.addWidget(file_frame)

        # Font Selection Frame
        font_frame = QFrame()
        font_layout = QHBoxLayout(font_frame)
        font_layout.setContentsMargins(10, 10, 10, 10)

        font_label = QLabel("Target Font:")
        font_label.setFixedWidth(80)
        font_layout.addWidget(font_label)

        self.font_entry = QLineEdit(self.font_name)
        font_layout.addWidget(self.font_entry)

        layout.addWidget(font_frame)

        # Action Frame
        action_frame = QFrame()
        action_layout = QVBoxLayout(action_frame)
        action_layout.setContentsMargins(10, 20, 10, 20)

        self.start_button = QPushButton("Start Processing")
        self.start_button.setFixedSize(150, 40)
        self.start_button.clicked.connect(self.process_file)
        action_layout.addWidget(self.start_button, alignment=Qt.AlignmentFlag.AlignCenter)

        layout.addWidget(action_frame)

        # Status Label
        self.status_label = QLabel("")
        self.status_label.setStyleSheet("color: green;")
        layout.addWidget(self.status_label, alignment=Qt.AlignmentFlag.AlignCenter)

    def browse_file(self):
        file_dialog = QFileDialog()
        file_dialog.setNameFilters([
            "Office Files (*.docx *.xlsx *.pptx)",
            "Word Documents (*.docx)",
            "Excel Workbooks (*.xlsx)",
            "PowerPoint Presentations (*.pptx)",
            "All files (*.*)"
        ])
        if file_dialog.exec():
            selected_files = file_dialog.selectedFiles()
            if selected_files:
                self.file_path = selected_files[0]
                self.file_entry.setText(self.file_path)
                self.status_label.setText("")
                self.status_label.setStyleSheet("color: green;")

    def process_file(self):
        path = self.file_path
        font = self.font_entry.text()

        if not path:
            QMessageBox.critical(self, "Error", "Please select a file first.")
            return
        if not font:
            QMessageBox.critical(self, "Error", "Please enter a target font name.")
            return

        self.status_label.setText("Processing...")
        self.status_label.setStyleSheet("color: blue;")
        QApplication.processEvents()

        try:
            file_dir, file_name = os.path.split(path)
            name, ext = os.path.splitext(file_name)
            output_path = os.path.join(file_dir, f"{name}_modified{ext}")

            if ext == ".docx":
                modified_doc = change_word_font(path, font)
                modified_doc.save(output_path)
            elif ext == ".xlsx":
                modified_workbook = change_excel_font(path, font)
                modified_workbook.save(output_path)
            elif ext == ".pptx":
                modified_prs = change_ppt_font(path, font)
                modified_prs.save(output_path)
            else:
                QMessageBox.critical(self, "Error", f"Unsupported file type: {ext}")
                self.status_label.setText("")
                self.status_label.setStyleSheet("color: red;")
                return

            self.status_label.setText(f"Success! Saved to {output_path}")
            self.status_label.setStyleSheet("color: green;")
            QMessageBox.information(self, "Success", "File processed successfully and saved as: " + output_path)

        except Exception as e:
            self.status_label.setText("An error occurred.")
            self.status_label.setStyleSheet("color: red;")
            QMessageBox.critical(self, "Error", "An error occurred during processing: " + str(e))


if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = FontUnifierApp()
    window.show()
    sys.exit(app.exec())
